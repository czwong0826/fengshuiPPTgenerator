import io
import uuid
import json

import streamlit as st
from PIL import Image
from streamlit_image_coordinates import streamlit_image_coordinates

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from knowledge_base import ROOM_KNOWLEDGE_BASE

st.set_page_config(page_title="风水报告生成器", layout="wide")

st.title("🏠 风水报告生成器（原型）")
st.caption(
    "按楼层来做：上传一张平面图，添加这一层的区域并打点、勾选内容，"
    "完成后点「➕ 添加新楼层」标签页继续下一层。全部做完后在页面底部生成 PPT / JSON。"
)

CUSTOM_TYPE = "自定义区域"
ROOM_TYPE_OPTIONS = list(ROOM_KNOWLEDGE_BASE.keys()) + [CUSTOM_TYPE]

# ---------- session_state 初始化 ----------
# floors：每个元素是一个"楼层"，楼层=一张图纸+挂在它上面的区域们，彼此完全独立：
#   {
#     "id": 楼层唯一ID,
#     "name": 楼层名称（如"底楼"）,
#     "image_bytes": 图片字节, "image_name": 文件名,
#     "image_width": 原始宽度, "image_height": 原始高度,
#     "instances": [ {"id":, "type":, "name":}, ... ],   # 这一层的区域实例
#     "coordinates": { instance_id: {"x":, "y":} },       # 该区域在这层图纸上的打点坐标
#     "add_area_form_version": 表单版本号（清空表单用）,
#   }
st.session_state.setdefault("floors", [])
st.session_state.setdefault("add_floor_form_version", 0)
st.session_state.setdefault("last_click_sig", {})  # 楼层ID -> 上一次点击签名，防止重复记录


def all_instances():
    """把所有楼层的区域实例拉平成一个列表，用于全局唯一命名计数。"""
    result = []
    for floor in st.session_state["floors"]:
        result.extend(floor["instances"])
    return result


def generate_default_name(room_type: str) -> str:
    """全局统一编号：不分楼层，第一个叫"卫生间"，第二个（不管在哪层）叫"卫生间 2"……
    这样方便用户在各处下拉框里凭名字分辨，不会重名混淆。
    """
    base = "未命名区域" if room_type == CUSTOM_TYPE else room_type
    count = sum(1 for inst in all_instances() if inst["type"] == room_type)
    if count == 0:
        return base
    return f"{base} {count + 1}"


# ---------- 渲染单个知识库条目（普通 checkbox 或 主选项+子选项），与原版逻辑一致 ----------
def render_knowledge_item(instance_id: str, idx: int, item):
    if isinstance(item, str):
        checked = st.checkbox(item, key=f"{instance_id}_checkbox_{idx}")
        return item if checked else None

    checked = st.checkbox(item["label"], key=f"{instance_id}_mainchk_{idx}")
    if not checked:
        return None

    values = {}
    all_filled = True
    with st.container():
        for group_key, group_cfg in item["choices"].items():
            options = group_cfg["options"]
            if group_cfg["multi"]:
                selected = st.multiselect(
                    "　↳ 请选择具体内容",
                    options=options,
                    key=f"{instance_id}_subms_{idx}_{group_key}",
                )
                if not selected:
                    all_filled = False
                else:
                    values[group_key] = "、".join(selected)
            else:
                selected = st.radio(
                    "　↳ 请选择具体内容",
                    options=options,
                    key=f"{instance_id}_subradio_{idx}_{group_key}",
                    index=None,
                )
                if selected is None:
                    all_filled = False
                else:
                    values[group_key] = selected

    if not all_filled:
        st.caption("⚠️ 请完成上方子选项的选择，否则该项不会计入报告")
        return None

    return item["template"].format(**values)


# ---------- 渲染单个楼层的完整编辑界面 ----------
def render_floor(floor: dict, report_data: dict):
    """渲染一个楼层的：删除入口 / 添加区域 / 已添加区域列表 / 打点标注 / 区域详情填写。
    填写结果写入 report_data（instance_id -> report），供后面生成 JSON / PPT 使用。
    """
    top_col1, top_col2 = st.columns([5, 1])
    top_col1.subheader(f"📐 {floor['name']}")
    if top_col2.button("🗑️ 删除该楼层", key=f"del_floor_{floor['id']}"):
        st.session_state["floors"] = [
            f for f in st.session_state["floors"] if f["id"] != floor["id"]
        ]
        st.session_state["pending_active_floor"] = None
        st.rerun()

    pil_image = Image.open(io.BytesIO(floor["image_bytes"]))
    img_width, img_height = floor["image_width"], floor["image_height"]

    # ---- 添加区域 ----
    # expanded 只在这个 expander 第一次出现时生效（默认展开），之后用户手动收起/展开的
    # 状态会通过 key 持久化，不会再被"添加了第一个区域"这种事情打断
    with st.expander("➕ 添加区域", expanded=True, key=f"add_area_expander_{floor['id']}"):
        form_version = floor["add_area_form_version"]
        with st.form(f"add_area_form_{floor['id']}_{form_version}", clear_on_submit=True):
            selected_type = st.selectbox(
                "房间类型", options=ROOM_TYPE_OPTIONS, key=f"new_type_{floor['id']}_{form_version}"
            )
            custom_name = st.text_input(
                "显示名称（可选，留空自动编号）",
                placeholder="例如：厕所 A",
                key=f"new_name_{floor['id']}_{form_version}",
            )
            submitted = st.form_submit_button("添加区域", use_container_width=True)
            if submitted:
                final_name = custom_name.strip() or generate_default_name(selected_type)
                floor["instances"].append(
                    {"id": uuid.uuid4().hex[:8], "type": selected_type, "name": final_name}
                )
                floor["add_area_form_version"] += 1
                st.toast(f"已在「{floor['name']}」添加「{final_name}」")
                st.rerun()

    # ---- 已添加区域列表 ----
    if floor["instances"]:
        with st.expander(f"📋 已添加区域（{len(floor['instances'])} 个）", expanded=False):
            for inst in list(floor["instances"]):
                c1, c2 = st.columns([4, 1])
                c1.write(f"📍 **{inst['name']}**　`{inst['type']}`")
                if c2.button("🗑️", key=f"del_inst_{inst['id']}", help="删除该区域"):
                    floor["instances"] = [i for i in floor["instances"] if i["id"] != inst["id"]]
                    floor["coordinates"].pop(inst["id"], None)
                    st.rerun()

    st.markdown("---")

    # ---- 打点标注 ----
    st.markdown("**🗺️ 打点标注**")
    st.caption(f"图纸「{floor['image_name']}」尺寸：{img_width} × {img_height} px。")

    if not floor["instances"]:
        # 还没有区域时，用同一个组件（streamlit_image_coordinates）显示图，
        # 保证跟"已有区域"时的显示尺寸规则完全一致，不会出现切换时图片忽大忽小；
        # 点击返回的坐标这里直接忽略，不记录、也不报错
        streamlit_image_coordinates(pil_image, key=f"coords_preview_{floor['id']}")
        st.info("请先在上方添加至少一个区域，才能开始为其记录坐标。")
    else:
        target_instance_id = st.selectbox(
            "选择要标注的区域（点击下方图片即可为该区域记录坐标）",
            options=[inst["id"] for inst in floor["instances"]],
            format_func=lambda iid: next(
                (i["name"] for i in floor["instances"] if i["id"] == iid), iid
            ),
            key=f"target_instance_{floor['id']}",
        )
        target_name = next(i["name"] for i in floor["instances"] if i["id"] == target_instance_id)
        st.caption(f"点击图片为「{target_name}」记录坐标。")

        coords = streamlit_image_coordinates(pil_image, key=f"coords_{floor['id']}")
        if coords is not None:
            sig = (coords["x"], coords["y"])
            if st.session_state["last_click_sig"].get(floor["id"]) != sig:
                st.session_state["last_click_sig"][floor["id"]] = sig
                floor["coordinates"][target_instance_id] = {"x": coords["x"], "y": coords["y"]}
                st.success(f"已记录「{target_name}」的坐标：({coords['x']}, {coords['y']})")

        if floor["coordinates"]:
            with st.expander("📌 本楼层已记录的坐标一览", expanded=False):
                readable = {
                    next((i["name"] for i in floor["instances"] if i["id"] == iid), iid): coord
                    for iid, coord in floor["coordinates"].items()
                }
                st.json(readable)
                if st.button("清空本楼层所有坐标", key=f"clear_coords_{floor['id']}"):
                    floor["coordinates"] = {}
                    st.session_state["last_click_sig"].pop(floor["id"], None)
                    st.rerun()

    st.markdown("---")

    # ---- 区域详情填写 ----
    if not floor["instances"]:
        return

    st.markdown("**📝 区域详情填写**")
    for inst in floor["instances"]:
        instance_id = inst["id"]
        area_type = inst["type"]
        area_name = inst["name"]

        with st.expander(f"📍 {area_name}", expanded=True):
            checked_items = []
            current_area_items = ROOM_KNOWLEDGE_BASE.get(area_type, [])
            if not current_area_items:
                st.caption("（该区域无预设勾选项，可在下方直接填写备注）")

            for idx, item in enumerate(current_area_items):
                result_text = render_knowledge_item(instance_id, idx, item)
                if result_text:
                    checked_items.append(result_text)

            note = st.text_area(
                "额外备注",
                key=f"{instance_id}_note",
                placeholder=f"请输入针对「{area_name}」的额外说明...",
                height=100,
            )

            report_data[instance_id] = {
                "名称": area_name,
                "类型": area_type,
                "楼层": floor["name"],
                "已勾选项": checked_items,
                "备注": note,
            }

            if instance_id in floor["coordinates"]:
                coord = floor["coordinates"][instance_id]
                report_data[instance_id]["坐标"] = coord
                st.caption(f"📌 已标注坐标：({coord['x']}, {coord['y']})")


# ---------- 1. 楼层标签页：已有楼层 + "添加新楼层" ----------
# 注意：st.tabs 用了显式 key，Streamlit 规定"带 key 的组件一旦实例化，
# 本次运行内后面的代码就不能再改它的 session_state"。所以"创建楼层后跳转到哪"
# 不能在下面表单提交处直接赋值，只能先记到 pending_active_floor 这个中转变量里，
# rerun 之后、st.tabs() 真正被调用之前，在这里统一写进它自己的 key
if "pending_active_floor" in st.session_state:
    st.session_state["floor_tabs"] = st.session_state.pop("pending_active_floor")

floors = st.session_state["floors"]
tab_labels = [f["name"] for f in floors] + ["➕ 添加新楼层"]
tabs = st.tabs(tab_labels, key="floor_tabs", on_change="rerun")

report_data = {}

for i, floor in enumerate(floors):
    with tabs[i]:
        render_floor(floor, report_data)

with tabs[-1]:
    st.subheader("添加新楼层")
    st.caption("每个楼层对应一张平面图，之后这一层的所有区域打点、生成 PPT 都只会用这张图。")
    form_version = st.session_state["add_floor_form_version"]
    with st.form(f"add_floor_form_{form_version}", clear_on_submit=True):
        floor_name = st.text_input(
            "楼层 / 页面名称", placeholder="例如：底楼、二楼", key=f"new_floor_name_{form_version}"
        )
        floor_image = st.file_uploader(
            "上传该楼层的平面图", type=["png", "jpg", "jpeg"], key=f"new_floor_image_{form_version}"
        )
        submitted = st.form_submit_button("创建楼层", use_container_width=True)
        if submitted:
            if not floor_name.strip():
                st.warning("请填写楼层名称。")
            elif floor_image is None:
                st.warning("请上传该楼层的平面图。")
            else:
                image_bytes = floor_image.getvalue()
                pil_img = Image.open(io.BytesIO(image_bytes))
                st.session_state["floors"].append(
                    {
                        "id": uuid.uuid4().hex[:8],
                        "name": floor_name.strip(),
                        "image_bytes": image_bytes,
                        "image_name": floor_image.name,
                        "image_width": pil_img.width,
                        "image_height": pil_img.height,
                        "instances": [],
                        "coordinates": {},
                        "add_area_form_version": 0,
                    }
                )
                st.session_state["add_floor_form_version"] += 1
                st.session_state["pending_active_floor"] = floor_name.strip()
                st.toast(f"已创建楼层「{floor_name.strip()}」")
                st.rerun()

st.markdown("---")


# ---------- PPT 生成相关工具函数 ----------
def add_arrowhead(connector, end="tail", arrow_type="triangle"):
    """python-pptx 没有现成的箭头 API，手动往连接线的 <a:ln> 里插入箭头描述。"""
    ln = connector.line._get_or_add_ln()
    tag = qn("a:tailEnd") if end == "tail" else qn("a:headEnd")
    for el in ln.findall(tag):
        ln.remove(el)
    el = ln.makeelement(tag, {"type": arrow_type, "w": "med", "len": "med"})
    ln.append(el)


def draw_marker(slide, center_x, center_y):
    """无填充红边小正方形 + 中心极小红点。"""
    square_size = Pt(14)
    square = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(center_x - square_size / 2),
        int(center_y - square_size / 2),
        square_size,
        square_size,
    )
    square.fill.background()
    square.line.color.rgb = RGBColor(255, 0, 0)
    square.line.width = Pt(1.5)
    square.shadow.inherit = False

    dot_size = Pt(3)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(center_x - dot_size / 2),
        int(center_y - dot_size / 2),
        dot_size,
        dot_size,
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor(255, 0, 0)
    dot.line.fill.background()
    dot.shadow.inherit = False


def add_callout(slide, left, top, width, height, title, bullets):
    """小号文字气泡：标题加粗下划线，内容"-"短横杠列点。
    框的高度是外面按"这一侧一共几个区域"平均分配好固定传进来的，
    这里打开 PowerPoint 原生的"自动缩小字体以适应形状"（TEXT_TO_FIT_SHAPE）：
    内容一旦超出框高，PowerPoint 打开时会自动把字号按比例缩小，不会出现文字溢出框外。
    （这个自动缩小是 PowerPoint 渲染时生效的，用 LibreOffice 等工具预览可能看不出缩小效果。）
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    tf.text = title
    title_p = tf.paragraphs[0]
    title_p.font.bold = True
    title_p.font.underline = True
    title_p.font.size = Pt(11)

    lines = list(bullets) if bullets else ["（未填写具体内容）"]
    for line in lines:
        p = tf.add_paragraph()
        p.text = f"- {line}"
        p.font.size = Pt(9.5)

    return box


def split_left_right(chunk_instances, coords_lookup, img_width, img_height):
    """按标记点在图上的左右半边分配气泡，并在数量差距过大时做一次平衡。
    没有坐标的区域按 0.5（中线）处理，参与平衡时优先被移动。
    返回 (left_list, right_list)，每个元素是 (inst, x_ratio, y_ratio_or_None)。
    """
    items = []
    for inst in chunk_instances:
        coord = coords_lookup.get(inst["id"])
        if coord:
            x_ratio = coord["x"] / img_width
            y_ratio = coord["y"] / img_height
        else:
            x_ratio, y_ratio = 0.5, None
        items.append((inst, x_ratio, y_ratio))

    left = [it for it in items if it[1] < 0.5]
    right = [it for it in items if it[1] >= 0.5]

    def move_closest_to_center(src, dst):
        src.sort(key=lambda it: abs(it[1] - 0.5))
        dst.append(src.pop(0))

    while len(left) - len(right) > 1:
        move_closest_to_center(left, right)
    while len(right) - len(left) > 1:
        move_closest_to_center(right, left)

    # 有坐标的按纵向位置从上到下排；没坐标的（y_ratio=None）统一排在最后
    left.sort(key=lambda it: (it[2] is None, it[2] if it[2] is not None else 0))
    right.sort(key=lambda it: (it[2] is None, it[2] if it[2] is not None else 0))
    return left, right


def chunk_list(lst, size):
    for i in range(0, len(lst), size):
        yield lst[i : i + size]


def build_pptx(floors: list, report_data: dict, max_per_page: int) -> bytes:
    """按楼层出页：同一楼层的区域合并展示，图纸居中，每个区域变成一个小气泡
    贴在图的左右两侧，用红色虚线肘形箭头指向它自己在图上的标记点。
    单个楼层区域数超过 max_per_page 时自动拆成"楼层名 (1/2)"这样的续页。
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_only_layout = prs.slide_layouts[5]

    CONTENT_TOP = Inches(1.3)
    CONTENT_BOTTOM = Inches(7.0)
    CALLOUT_WIDTH = Inches(2.6)
    CALLOUT_GAP = Inches(0.3)
    SIDE_MARGIN = Inches(0.3)
    CALLOUT_VGAP = Inches(0.15)

    left_zone_right_edge = SIDE_MARGIN + CALLOUT_WIDTH
    right_zone_left_edge = prs.slide_width - SIDE_MARGIN - CALLOUT_WIDTH

    img_max_left = left_zone_right_edge + CALLOUT_GAP
    img_max_right = right_zone_left_edge - CALLOUT_GAP
    img_max_width = img_max_right - img_max_left
    img_max_height = CONTENT_BOTTOM - CONTENT_TOP

    for floor in floors:
        instances = floor["instances"]
        if not instances:
            continue

        chunks = list(chunk_list(instances, max_per_page))
        total_pages = len(chunks)

        for page_idx, chunk in enumerate(chunks, start=1):
            slide = prs.slides.add_slide(title_only_layout)
            title_text = floor["name"] if total_pages == 1 else f"{floor['name']} ({page_idx}/{total_pages})"
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_shape.left = Inches(0.4)
            title_shape.top = Inches(0.25)
            title_shape.width = prs.slide_width - Inches(0.8)
            title_shape.height = Inches(0.8)
            title_paragraph = title_shape.text_frame.paragraphs[0]
            title_paragraph.font.size = Pt(28)
            title_paragraph.alignment = PP_ALIGN.LEFT

            # ---- 居中放置该楼层的图纸 ----
            aspect = floor["image_width"] / floor["image_height"]
            box_aspect = img_max_width / img_max_height
            if aspect > box_aspect:
                final_width = img_max_width
                final_height = int(final_width / aspect)
            else:
                final_height = img_max_height
                final_width = int(final_height * aspect)

            pic_left = img_max_left + int((img_max_width - final_width) / 2)
            pic_top = CONTENT_TOP + int((img_max_height - final_height) / 2)

            picture = slide.shapes.add_picture(
                io.BytesIO(floor["image_bytes"]),
                left=pic_left,
                top=pic_top,
                width=final_width,
                height=final_height,
            )

            # ---- 分配左右两侧气泡 ----
            left_items, right_items = split_left_right(
                chunk, floor["coordinates"], floor["image_width"], floor["image_height"]
            )

            def render_side(items, side_left_x, anchor="left"):
                if not items:
                    return
                # 把这一侧的可用纵向空间，按区域个数平均分配成等高的框——
                # 不再靠字符数瞎猜高度，配合 add_callout 里打开的自动缩字，
                # 内容多的框会自动缩小字号塞进去，不会溢出、也不会互相重叠
                n = len(items)
                available_height = CONTENT_BOTTOM - CONTENT_TOP
                box_height = int((available_height - CALLOUT_VGAP * (n - 1)) / n)

                current_top = CONTENT_TOP
                for inst, x_ratio, y_ratio in items:
                    report = report_data.get(inst["id"], {})
                    area_name = report.get("名称", inst["name"])
                    checked_items = report.get("已勾选项", [])
                    note = report.get("备注", "").strip()
                    bullets = list(checked_items)
                    if note:
                        # 备注按用户实际换行拆开，每一行单独变成一条"- "列点
                        bullets.extend(line.strip() for line in note.splitlines() if line.strip())

                    box = add_callout(
                        slide, side_left_x, current_top, CALLOUT_WIDTH, box_height, area_name, bullets
                    )

                    coord = floor["coordinates"].get(inst["id"])
                    if coord is not None:
                        marker_x = picture.left + int(picture.width * (coord["x"] / floor["image_width"]))
                        marker_y = picture.top + int(picture.height * (coord["y"] / floor["image_height"]))
                        draw_marker(slide, marker_x, marker_y)

                        begin_y = box.top + int(box.height / 2)
                        begin_x = (box.left + box.width) if anchor == "left" else box.left
                        connector = slide.shapes.add_connector(
                            MSO_CONNECTOR.ELBOW, begin_x, begin_y, marker_x, marker_y
                        )
                        connector.line.color.rgb = RGBColor(255, 0, 0)
                        connector.line.width = Pt(1)
                        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                        add_arrowhead(connector, end="tail", arrow_type="triangle")

                    current_top = current_top + box_height + CALLOUT_VGAP

            render_side(left_items, SIDE_MARGIN, anchor="left")
            render_side(right_items, right_zone_left_edge, anchor="right")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ---------- 2. 生成 JSON 数据 / PPT 文件 ----------
st.subheader("📤 生成报告")

max_per_page = st.number_input(
    "每页 PPT 最多显示多少个区域（超过自动拆成续页）",
    min_value=2,
    max_value=12,
    value=6,
    step=1,
)

ordered_reports = []
for floor in floors:
    for inst in floor["instances"]:
        if inst["id"] in report_data:
            ordered_reports.append(report_data[inst["id"]])

col1, col2 = st.columns(2)

with col1:
    if st.button("生成 JSON 数据", type="primary", use_container_width=True):
        if not floors:
            st.warning("尚未添加任何楼层，无法生成报告。")
        else:
            final_output = {
                "楼层列表": [
                    {
                        "楼层名称": floor["name"],
                        "区域列表": [
                            report_data[inst["id"]]
                            for inst in floor["instances"]
                            if inst["id"] in report_data
                        ],
                    }
                    for floor in floors
                ]
            }
            st.success("JSON 数据生成成功！")
            st.json(final_output)

            json_str = json.dumps(final_output, ensure_ascii=False, indent=2)
            st.download_button(
                label="下载 JSON 文件",
                data=json_str,
                file_name="fengshui_report.json",
                mime="application/json",
            )

with col2:
    if st.button("生成 PPT", type="secondary", use_container_width=True):
        if not floors:
            st.warning("尚未添加任何楼层，无法生成报告。")
        elif not any(floor["instances"] for floor in floors):
            st.warning("所有楼层都还没有添加区域，无法生成报告。")
        else:
            pptx_bytes = build_pptx(floors, report_data, int(max_per_page))
            st.success("PPT 生成成功！")
            st.download_button(
                label="下载 PPTX 文件",
                data=pptx_bytes,
                file_name="fengshui_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )